"""Tool definitions and executor for ProxyAPI tool calling."""

from __future__ import annotations

import fnmatch
import json
import os
import re
import subprocess
import urllib.error
import urllib.request
from pathlib import Path
from typing import TYPE_CHECKING

_NO_WINDOW = {"creationflags": subprocess.CREATE_NO_WINDOW} if os.name == "nt" else {}

from loguru import logger

from .email import EMAIL_TOOL_DEFINITIONS, execute_email_tool
from .calendar import CALENDAR_TOOL_DEFINITIONS, execute_calendar_tool

if TYPE_CHECKING:
    from agent_commander.session.extension_store import ExtensionStore
    from agent_commander.session.skill_store import SkillStore
    from agent_commander.session.project_store import ProjectStore

MAX_RESULT_CHARS = 32_000
COMMAND_TIMEOUT_S = 30
WEB_FETCH_TIMEOUT_S = 15

TOOL_DEFINITIONS = [
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "Read the contents of a file at the given path.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Absolute or relative file path to read.",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": "Write content to a file. Creates parent directories if needed.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Absolute or relative file path to write.",
                    },
                    "content": {
                        "type": "string",
                        "description": "Content to write to the file.",
                    },
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_directory",
            "description": "List files and directories at the given path.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Directory path to list. Defaults to current working directory.",
                    },
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_command",
            "description": "Execute a shell command and return its output.",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {
                        "type": "string",
                        "description": "Shell command to execute.",
                    },
                },
                "required": ["command"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "edit_file",
            "description": "Perform a surgical find-and-replace edit in a file. The old_string must match exactly (including whitespace and indentation). Use this instead of write_file when modifying existing files.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "File path to edit.",
                    },
                    "old_string": {
                        "type": "string",
                        "description": "Exact text to find in the file. Must be unique within the file.",
                    },
                    "new_string": {
                        "type": "string",
                        "description": "Text to replace old_string with.",
                    },
                },
                "required": ["path", "old_string", "new_string"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "glob",
            "description": "Find files matching a glob pattern (e.g., '**/*.py', 'src/**/*.ts'). Returns matching file paths.",
            "parameters": {
                "type": "object",
                "properties": {
                    "pattern": {
                        "type": "string",
                        "description": "Glob pattern to match files (e.g., '**/*.py').",
                    },
                    "path": {
                        "type": "string",
                        "description": "Base directory to search in (defaults to cwd).",
                    },
                },
                "required": ["pattern"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "grep",
            "description": "Search file contents for a regex pattern. Returns matching lines with file paths and line numbers.",
            "parameters": {
                "type": "object",
                "properties": {
                    "pattern": {
                        "type": "string",
                        "description": "Regex pattern to search for.",
                    },
                    "path": {
                        "type": "string",
                        "description": "File or directory to search in (defaults to cwd).",
                    },
                    "include": {
                        "type": "string",
                        "description": "Glob pattern to filter files (e.g., '*.py', '*.ts').",
                    },
                },
                "required": ["pattern"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_fetch",
            "description": "Fetch content from a URL. Returns the response body as text.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {
                        "type": "string",
                        "description": "URL to fetch.",
                    },
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_document",
            "description": "Read a Word document (.docx) and return its text content. Preserves paragraph structure.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the .docx file.",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_document",
            "description": "Create a Word document (.docx) from text. Lines starting with '# ' become headings. Creates parent directories if needed.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path for the output .docx file.",
                    },
                    "content": {
                        "type": "string",
                        "description": "Text content. Lines starting with '# ' become Heading 1, '## ' become Heading 2, '### ' become Heading 3. Other lines become normal paragraphs.",
                    },
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_spreadsheet",
            "description": "Read an Excel spreadsheet (.xlsx) and return cell data as text. By default reads the active sheet.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the .xlsx file.",
                    },
                    "sheet": {
                        "type": "string",
                        "description": "Sheet name to read (optional, defaults to active sheet).",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_spreadsheet",
            "description": "Create or overwrite an Excel spreadsheet (.xlsx). Accepts data as a JSON array of rows (each row is an array of cell values). First row is typically headers.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path for the output .xlsx file.",
                    },
                    "data": {
                        "type": "array",
                        "items": {
                            "type": "array",
                            "items": {
                                "anyOf": [
                                    {"type": "string"},
                                    {"type": "number"},
                                    {"type": "boolean"},
                                    {"type": "null"},
                                ]
                            },
                        },
                        "description": "Array of rows. Each row is an array of cell values. Example: [[\"Name\", \"Age\"], [\"Alice\", 30], [\"Bob\", 25]]",
                    },
                    "sheet": {
                        "type": "string",
                        "description": "Sheet name (optional, defaults to 'Sheet').",
                    },
                },
                "required": ["path", "data"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_presentation",
            "description": "Read a PowerPoint presentation (.pptx) and return slide text content.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path to the .pptx file.",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_presentation",
            "description": "Create a PowerPoint presentation (.pptx). Each slide has a title and optional body text.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Path for the output .pptx file.",
                    },
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "body": {"type": "string"},
                            },
                        },
                        "description": "Array of slide objects. Each has 'title' (required) and 'body' (optional, supports newlines for bullet points).",
                    },
                },
                "required": ["path", "slides"],
            },
        },
    },
    *EMAIL_TOOL_DEFINITIONS,
    *CALENDAR_TOOL_DEFINITIONS,
]

# ---------------------------------------------------------------------------
# Team & Project tool definitions (appended when stores are available)
# ---------------------------------------------------------------------------

TEAM_PROJECT_TOOL_DEFINITIONS = [
    {
        "type": "function",
        "function": {
            "name": "team_list_roles",
            "description": "List all roles in the Team skill library. Returns names, categories and descriptions.",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "team_upsert_role",
            "description": (
                "Create a new role or update an existing one (matched by name) in the Team skill library. "
                "Use this to capture reusable instructions, personas, or expertise profiles."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "Role name (unique identifier)."},
                    "category": {"type": "string", "description": "Category tag, e.g. 'Dev', 'Design', 'Management'."},
                    "description": {"type": "string", "description": "Short one-line description."},
                    "content": {"type": "string", "description": "Full role instructions in Markdown."},
                },
                "required": ["name", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "project_list",
            "description": "List all projects with their checklist progress.",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "project_upsert",
            "description": (
                "Create a new project or update an existing one (matched by name). "
                "Use this to track goals, set the project path, and maintain a checklist of milestones."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "Project name (unique identifier)."},
                    "description": {"type": "string", "description": "Project description or goal."},
                    "workdir": {"type": "string", "description": "Absolute path to the project directory (optional)."},
                    "checklist": {
                        "type": "array",
                        "description": "Full checklist for the project. Replaces existing checklist.",
                        "items": {
                            "type": "object",
                            "properties": {
                                "text": {"type": "string"},
                                "done": {"type": "boolean"},
                            },
                            "required": ["text"],
                        },
                    },
                },
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "project_set_checklist_item",
            "description": (
                "Mark a specific checklist item as done or not done in a project. "
                "Use this to track progress as you complete milestones."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "project_name": {"type": "string", "description": "Name of the project."},
                    "item_text": {"type": "string", "description": "Exact text of the checklist item to update."},
                    "done": {"type": "boolean", "description": "True to mark as done, false to unmark."},
                },
                "required": ["project_name", "item_text", "done"],
            },
        },
    },
]


def execute_tool(
    name: str,
    arguments_json: str,
    cwd: str | None = None,
    extension_store: "ExtensionStore | None" = None,
    active_extension_ids: list[str] | None = None,
    skill_store: "SkillStore | None" = None,
    project_store: "ProjectStore | None" = None,
) -> str:
    """Execute a tool by name and return the result as a string."""
    try:
        args = json.loads(arguments_json) if isinstance(arguments_json, str) else arguments_json
    except json.JSONDecodeError as exc:
        return f"Error: invalid tool arguments JSON: {exc}"

    if not isinstance(args, dict):
        return "Error: tool arguments must be a JSON object"

    try:
        if name == "read_file":
            return _read_file(args.get("path", ""), cwd)
        elif name == "write_file":
            return _write_file(args.get("path", ""), args.get("content", ""), cwd)
        elif name == "edit_file":
            return _edit_file(args.get("path", ""), args.get("old_string", ""), args.get("new_string", ""), cwd)
        elif name == "list_directory":
            return _list_directory(args.get("path", "."), cwd)
        elif name == "glob":
            return _glob(args.get("pattern", ""), args.get("path", ""), cwd)
        elif name == "grep":
            return _grep(args.get("pattern", ""), args.get("path", ""), args.get("include", ""), cwd)
        elif name == "run_command":
            return _run_command(args.get("command", ""), cwd)
        elif name == "web_fetch":
            return _web_fetch(args.get("url", ""))
        elif name == "read_document":
            return _read_document(args.get("path", ""), cwd)
        elif name == "write_document":
            return _write_document(args.get("path", ""), args.get("content", ""), cwd)
        elif name == "read_spreadsheet":
            return _read_spreadsheet(args.get("path", ""), args.get("sheet", ""), cwd)
        elif name == "write_spreadsheet":
            return _write_spreadsheet(args.get("path", ""), args.get("data", []), args.get("sheet", ""), cwd)
        elif name == "read_presentation":
            return _read_presentation(args.get("path", ""), cwd)
        elif name == "write_presentation":
            return _write_presentation(args.get("path", ""), args.get("slides", []), cwd)
        elif name.startswith("email_"):
            if extension_store is None:
                return "Error: email tools require extensions (no extension_store provided)"
            return execute_email_tool(
                name=name, args=args, extension_store=extension_store,
                active_ids=active_extension_ids,
            )
        elif name.startswith("calendar_"):
            if extension_store is None:
                return "Error: calendar tools require extensions (no extension_store provided)"
            return execute_calendar_tool(
                name=name, args=args, extension_store=extension_store,
                active_ids=active_extension_ids,
            )
        elif name == "team_list_roles":
            return _team_list_roles(skill_store)
        elif name == "team_upsert_role":
            return _team_upsert_role(
                args.get("name", ""), args.get("category", ""),
                args.get("description", ""), args.get("content", ""), skill_store,
            )
        elif name == "project_list":
            return _project_list(project_store)
        elif name == "project_upsert":
            return _project_upsert(
                args.get("name", ""), args.get("description", ""),
                args.get("workdir", ""), args.get("checklist", []), project_store,
            )
        elif name == "project_set_checklist_item":
            return _project_set_checklist_item(
                args.get("project_name", ""), args.get("item_text", ""),
                bool(args.get("done", False)), project_store,
            )
        else:
            return f"Error: unknown tool '{name}'"
    except Exception as exc:
        logger.warning(f"Tool '{name}' failed: {exc}")
        return f"Error executing {name}: {exc}"


# ---------------------------------------------------------------------------
# Team & Project tool implementations
# ---------------------------------------------------------------------------


def _team_list_roles(skill_store: "SkillStore | None") -> str:
    if skill_store is None:
        return "Error: team skill store not available"
    skills = skill_store.list_skills()
    if not skills:
        return "No roles defined yet."
    lines = ["Roles in Team library:\n"]
    for s in skills:
        cat = f"  [{s.category}]" if s.category else ""
        lines.append(f"- **{s.name}**{cat}: {s.description}")
    return "\n".join(lines)


def _team_upsert_role(
    name: str, category: str, description: str, content: str,
    skill_store: "SkillStore | None",
) -> str:
    if skill_store is None:
        return "Error: team skill store not available"
    if not name:
        return "Error: name is required"
    # Find existing by name
    existing = next((s for s in skill_store.list_skills() if s.name.lower() == name.lower()), None)
    if existing:
        skill_store.update_skill(existing.id, name, description, category, content)
        return f"Role '{name}' updated successfully."
    else:
        skill_store.create_skill(name, description, category, content)
        return f"Role '{name}' created successfully."


def _project_list(project_store: "ProjectStore | None") -> str:
    if project_store is None:
        return "Error: project store not available"
    projects = project_store.list_projects()
    if not projects:
        return "No projects defined yet."
    lines = ["Projects:\n"]
    for p in projects:
        done = sum(1 for item in p.checklist if item.get("done"))
        total = len(p.checklist)
        progress = f" [{done}/{total} done]" if total else ""
        path = f" — {p.workdir}" if p.workdir else ""
        lines.append(f"- **{p.name}**{progress}{path}: {p.description}")
    return "\n".join(lines)


def _project_upsert(
    name: str, description: str, workdir: str, checklist: list,
    project_store: "ProjectStore | None",
) -> str:
    if project_store is None:
        return "Error: project store not available"
    if not name:
        return "Error: name is required"
    # Normalise checklist items
    normalised = [
        {"text": str(item.get("text", "")), "done": bool(item.get("done", False))}
        for item in checklist if item.get("text")
    ]
    existing = next((p for p in project_store.list_projects() if p.name.lower() == name.lower()), None)
    if existing:
        existing.description = description or existing.description
        if workdir:
            existing.workdir = workdir
        if normalised:
            existing.checklist = normalised
        project_store.update_project(existing)
        return f"Project '{name}' updated successfully."
    else:
        meta = project_store.create_project(name, description, workdir)
        if normalised:
            meta.checklist = normalised
            project_store.update_project(meta)
        return f"Project '{name}' created successfully."


def _project_set_checklist_item(
    project_name: str, item_text: str, done: bool,
    project_store: "ProjectStore | None",
) -> str:
    if project_store is None:
        return "Error: project store not available"
    if not project_name or not item_text:
        return "Error: project_name and item_text are required"
    project = next(
        (p for p in project_store.list_projects() if p.name.lower() == project_name.lower()),
        None,
    )
    if project is None:
        return f"Error: project '{project_name}' not found"
    matched = False
    for item in project.checklist:
        if item.get("text", "").lower() == item_text.lower():
            item["done"] = done
            matched = True
            break
    if not matched:
        # Item not found — add it
        project.checklist.append({"text": item_text, "done": done})
    project_store.update_project(project)
    status = "done" if done else "not done"
    action = "marked" if matched else "added and marked"
    return f"Checklist item '{item_text}' {action} as {status} in project '{project_name}'."


def _resolve_path(raw_path: str, cwd: str | None) -> Path:
    """Resolve a path, using cwd as base for relative paths."""
    p = Path(raw_path).expanduser()
    if not p.is_absolute() and cwd:
        p = Path(cwd) / p
    return p.resolve()


def _truncate(text: str) -> str:
    if len(text) > MAX_RESULT_CHARS:
        return text[:MAX_RESULT_CHARS] + f"\n... (truncated, {len(text)} total chars)"
    return text


def _read_file(path: str, cwd: str | None) -> str:
    if not path:
        return "Error: path is required"
    resolved = _resolve_path(path, cwd)
    if not resolved.is_file():
        return f"Error: file not found: {resolved}"
    content = resolved.read_text(encoding="utf-8", errors="replace")
    logger.info(f"Tool read_file: {resolved} ({len(content)} chars)")
    return _truncate(content)


def _write_file(path: str, content: str, cwd: str | None) -> str:
    if not path:
        return "Error: path is required"
    resolved = _resolve_path(path, cwd)
    resolved.parent.mkdir(parents=True, exist_ok=True)
    resolved.write_text(content, encoding="utf-8")
    logger.info(f"Tool write_file: {resolved} ({len(content)} chars)")
    return f"File written: {resolved} ({len(content)} chars)"


def _list_directory(path: str, cwd: str | None) -> str:
    resolved = _resolve_path(path or ".", cwd)
    if not resolved.is_dir():
        return f"Error: directory not found: {resolved}"
    entries: list[str] = []
    try:
        for item in sorted(resolved.iterdir()):
            kind = "dir" if item.is_dir() else "file"
            size = ""
            if item.is_file():
                try:
                    size = f" ({item.stat().st_size} bytes)"
                except OSError:
                    pass
            entries.append(f"  [{kind}] {item.name}{size}")
    except PermissionError:
        return f"Error: permission denied: {resolved}"
    header = f"Directory: {resolved}\n"
    if not entries:
        return header + "  (empty)"
    return _truncate(header + "\n".join(entries))


def _run_command(command: str, cwd: str | None) -> str:
    if not command:
        return "Error: command is required"
    work_dir = cwd or os.getcwd()
    logger.info(f"Tool run_command: {command!r} in {work_dir}")
    try:
        result = subprocess.run(
            command,
            shell=True,
            cwd=work_dir,
            capture_output=True,
            text=True,
            timeout=COMMAND_TIMEOUT_S,
            **_NO_WINDOW,
        )
        parts: list[str] = []
        if result.stdout:
            parts.append(result.stdout)
        if result.stderr:
            parts.append(f"STDERR:\n{result.stderr}")
        if result.returncode != 0:
            parts.append(f"Exit code: {result.returncode}")
        output = "\n".join(parts) if parts else "(no output)"
        return _truncate(output)
    except subprocess.TimeoutExpired:
        return f"Error: command timed out after {COMMAND_TIMEOUT_S}s"


def _edit_file(path: str, old_string: str, new_string: str, cwd: str | None) -> str:
    if not path:
        return "Error: path is required"
    if not old_string:
        return "Error: old_string is required"
    resolved = _resolve_path(path, cwd)
    if not resolved.is_file():
        return f"Error: file not found: {resolved}"
    content = resolved.read_text(encoding="utf-8", errors="replace")
    count = content.count(old_string)
    if count == 0:
        return f"Error: old_string not found in {resolved}"
    if count > 1:
        return f"Error: old_string found {count} times in {resolved} — must be unique. Provide more surrounding context."
    new_content = content.replace(old_string, new_string, 1)
    resolved.write_text(new_content, encoding="utf-8")
    logger.info(f"Tool edit_file: {resolved} (replaced 1 occurrence)")
    return f"File edited: {resolved} (1 replacement made)"


def _glob(pattern: str, path: str, cwd: str | None) -> str:
    if not pattern:
        return "Error: pattern is required"
    base = _resolve_path(path or ".", cwd)
    if not base.is_dir():
        return f"Error: directory not found: {base}"
    matches: list[str] = []
    max_results = 500
    for match in base.glob(pattern):
        matches.append(str(match))
        if len(matches) >= max_results:
            break
    if not matches:
        return f"No files matching '{pattern}' in {base}"
    header = f"Found {len(matches)} file(s) matching '{pattern}' in {base}:\n"
    return _truncate(header + "\n".join(matches))


def _grep(pattern: str, path: str, include: str, cwd: str | None) -> str:
    if not pattern:
        return "Error: pattern is required"
    base = _resolve_path(path or ".", cwd)
    try:
        regex = re.compile(pattern, re.IGNORECASE)
    except re.error as exc:
        return f"Error: invalid regex: {exc}"

    matches: list[str] = []
    max_matches = 200
    files_searched = 0

    if base.is_file():
        file_list = [base]
    else:
        file_list = sorted(base.rglob("*"))

    for file_path in file_list:
        if not file_path.is_file():
            continue
        if include and not fnmatch.fnmatch(file_path.name, include):
            continue
        # Skip binary/large files
        try:
            size = file_path.stat().st_size
            if size > 1_000_000:
                continue
        except OSError:
            continue
        try:
            text = file_path.read_text(encoding="utf-8", errors="ignore")
        except (PermissionError, OSError):
            continue
        files_searched += 1
        for line_num, line in enumerate(text.splitlines(), 1):
            if regex.search(line):
                rel = file_path.relative_to(base) if base.is_dir() and file_path.is_relative_to(base) else file_path
                matches.append(f"{rel}:{line_num}: {line.rstrip()}")
                if len(matches) >= max_matches:
                    break
        if len(matches) >= max_matches:
            break

    if not matches:
        return f"No matches for '{pattern}' in {base} ({files_searched} files searched)"
    header = f"Found {len(matches)} match(es) for '{pattern}' ({files_searched} files searched):\n"
    return _truncate(header + "\n".join(matches))


def _try_import_docx():
    try:
        import docx
        return docx
    except ImportError:
        return None


def _try_import_openpyxl():
    try:
        import openpyxl
        return openpyxl
    except ImportError:
        return None


def _try_import_pptx():
    try:
        import pptx
        return pptx
    except ImportError:
        return None


def _read_document(path: str, cwd: str | None) -> str:
    if not path:
        return "Error: path is required"
    docx = _try_import_docx()
    if docx is None:
        return "Error: python-docx is not installed (pip install python-docx)"
    resolved = _resolve_path(path, cwd)
    if not resolved.is_file():
        return f"Error: file not found: {resolved}"
    try:
        doc = docx.Document(str(resolved))
    except Exception as exc:
        return f"Error: cannot open document: {exc}"
    parts: list[str] = []
    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        text = para.text
        if "heading 1" in style:
            parts.append(f"# {text}")
        elif "heading 2" in style:
            parts.append(f"## {text}")
        elif "heading 3" in style:
            parts.append(f"### {text}")
        else:
            parts.append(text)
    # Also extract table content
    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("")
            parts.append(" | ".join(["---"] * len(table.rows[0].cells)))
            parts.extend(rows)
    content = "\n".join(parts)
    logger.info(f"Tool read_document: {resolved} ({len(content)} chars)")
    return _truncate(content)


def _write_document(path: str, content: str, cwd: str | None) -> str:
    if not path:
        return "Error: path is required"
    docx = _try_import_docx()
    if docx is None:
        return "Error: python-docx is not installed (pip install python-docx)"
    resolved = _resolve_path(path, cwd)
    resolved.parent.mkdir(parents=True, exist_ok=True)
    doc = docx.Document()
    for line in content.split("\n"):
        if line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        else:
            doc.add_paragraph(line)
    doc.save(str(resolved))
    logger.info(f"Tool write_document: {resolved}")
    return f"Document written: {resolved}"


def _read_spreadsheet(path: str, sheet: str, cwd: str | None) -> str:
    if not path:
        return "Error: path is required"
    openpyxl = _try_import_openpyxl()
    if openpyxl is None:
        return "Error: openpyxl is not installed (pip install openpyxl)"
    resolved = _resolve_path(path, cwd)
    if not resolved.is_file():
        return f"Error: file not found: {resolved}"
    try:
        wb = openpyxl.load_workbook(str(resolved), read_only=True, data_only=True)
    except Exception as exc:
        return f"Error: cannot open spreadsheet: {exc}"
    try:
        if sheet:
            if sheet not in wb.sheetnames:
                wb.close()
                return f"Error: sheet '{sheet}' not found. Available: {', '.join(wb.sheetnames)}"
            ws = wb[sheet]
        else:
            ws = wb.active
        sheet_name = ws.title if ws else "?"
        lines: list[str] = [f"Sheet: {sheet_name}", f"Available sheets: {', '.join(wb.sheetnames)}", ""]
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("\t".join(cells))
            row_count += 1
            if row_count >= 5000:
                lines.append(f"... (truncated at {row_count} rows)")
                break
    finally:
        wb.close()
    content = "\n".join(lines)
    logger.info(f"Tool read_spreadsheet: {resolved} ({row_count} rows)")
    return _truncate(content)


def _write_spreadsheet(path: str, data: list, sheet: str, cwd: str | None) -> str:
    if not path:
        return "Error: path is required"
    if not data:
        return "Error: data is required (array of rows)"
    openpyxl = _try_import_openpyxl()
    if openpyxl is None:
        return "Error: openpyxl is not installed (pip install openpyxl)"
    resolved = _resolve_path(path, cwd)
    resolved.parent.mkdir(parents=True, exist_ok=True)
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet or "Sheet"
    for row in data:
        if isinstance(row, list):
            ws.append(row)
        else:
            ws.append([row])
    wb.save(str(resolved))
    wb.close()
    logger.info(f"Tool write_spreadsheet: {resolved} ({len(data)} rows)")
    return f"Spreadsheet written: {resolved} ({len(data)} rows)"


def _read_presentation(path: str, cwd: str | None) -> str:
    if not path:
        return "Error: path is required"
    pptx = _try_import_pptx()
    if pptx is None:
        return "Error: python-pptx is not installed (pip install python-pptx)"
    resolved = _resolve_path(path, cwd)
    if not resolved.is_file():
        return f"Error: file not found: {resolved}"
    try:
        prs = pptx.Presentation(str(resolved))
    except Exception as exc:
        return f"Error: cannot open presentation: {exc}"
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))
        parts.append("")
    content = "\n".join(parts)
    logger.info(f"Tool read_presentation: {resolved} ({len(prs.slides)} slides)")
    return _truncate(content)


def _write_presentation(path: str, slides: list, cwd: str | None) -> str:
    if not path:
        return "Error: path is required"
    if not slides:
        return "Error: slides is required (array of slide objects)"
    pptx = _try_import_pptx()
    if pptx is None:
        return "Error: python-pptx is not installed (pip install python-pptx)"
    resolved = _resolve_path(path, cwd)
    resolved.parent.mkdir(parents=True, exist_ok=True)
    prs = pptx.Presentation()
    for slide_data in slides:
        if not isinstance(slide_data, dict):
            continue
        title = slide_data.get("title", "")
        body = slide_data.get("body", "")
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        if body and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            lines = body.split("\n")
            for j, line in enumerate(lines):
                if j == 0:
                    tf.paragraphs[0].text = line
                else:
                    tf.add_paragraph().text = line
    prs.save(str(resolved))
    logger.info(f"Tool write_presentation: {resolved} ({len(slides)} slides)")
    return f"Presentation written: {resolved} ({len(slides)} slides)"


def _web_fetch(url: str) -> str:
    if not url:
        return "Error: url is required"
    if not url.startswith(("http://", "https://")):
        url = "https://" + url
    logger.info(f"Tool web_fetch: {url}")
    headers = {"User-Agent": "agent-commander-gui/1.0"}
    req = urllib.request.Request(url, headers=headers, method="GET")
    try:
        with urllib.request.urlopen(req, timeout=WEB_FETCH_TIMEOUT_S) as resp:
            content_type = (resp.headers.get("Content-Type") or "").lower()
            body = resp.read().decode("utf-8", errors="replace")
            # Strip HTML tags for readability if it's HTML
            if "html" in content_type:
                body = re.sub(r"<script[^>]*>.*?</script>", "", body, flags=re.DOTALL)
                body = re.sub(r"<style[^>]*>.*?</style>", "", body, flags=re.DOTALL)
                body = re.sub(r"<[^>]+>", " ", body)
                body = re.sub(r"\s+", " ", body).strip()
            return _truncate(body)
    except urllib.error.HTTPError as exc:
        return f"Error: HTTP {exc.code} for {url}"
    except urllib.error.URLError as exc:
        return f"Error: could not fetch {url}: {exc.reason}"
    except Exception as exc:
        return f"Error: {exc}"
